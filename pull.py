"""Fetch the LibCal feed and rebuild the SCDS upcoming-workshops slide.

Reads the live LibCal feed, groups events into the four program-area boxes,
and writes a dated copy of SCDS Slide.pptx with the event lines refreshed.
Formatting (fonts, colours, boxes, QR code) comes from the original deck.
"""
import copy
import sys
from datetime import datetime
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(errors="replace")

import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry

FEED_URL = ("https://libcal.mcmaster.ca/api_events.php"
            "?m=upc&cid=8132_7565&audience=&c=&d=&tags=&l=50&tar=0")
DECK_PATH = Path("SCDS Slide template.pptx")

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x0F, 0x0F, 0x0F)

# Priority order: an event tagged with several categories lands in the first
# box whose keywords match one of its LibCal category names (exact match).
BOXES = [
    {"header": "Data Analysis Support Hub", "color": WHITE,
     "keywords": ["dash", "data analysis support hub"]},
    {"header": "Research Data Management", "color": DARK,
     "keywords": ["research data management", "rdm"]},
    {"header": "Digital Research", "color": DARK,
     "keywords": ["digital research", "dr"]},
    {"header": "Do More with Digital Scholarship", "color": WHITE,
     "keywords": ["dmds"]},
]
BOX_HEADERS = {box["header"] for box in BOXES}

# Max *text lines* per section (titles wrap); keep in sync with docs/lib.js.
# DMDS shares its quadrant with an overflow textbox, hence 7.
MAX_LINES = {
    "Data Analysis Support Hub": 6,
    "Research Data Management": 8,
    "Digital Research": 6,
    "Do More with Digital Scholarship": 7,
}

# Approximate Aptos 16pt character capacity of the 5.67in-wide section
# boxes, calibrated so medium titles don't overcount wrapped lines. Counts
# the rendered "Month D - " prefix, not just the title.
CHARS_PER_LINE = 45

# Do More stacks a short main textbox over a narrow overflow textbox that
# stops short of the QR code: 2 lines up top, 5 below (7 total, matching
# MAX_LINES).
DMDS_MAIN_LINES = 2


def fetch_feed(url=FEED_URL, attempts=3):
    session = requests.Session()
    retry = Retry(total=attempts - 1, backoff_factor=1,
                  status_forcelist=(429, 500, 502, 503, 504))
    session.mount("https://", HTTPAdapter(max_retries=retry))
    response = session.get(url, timeout=30, headers={"User-Agent": "Mozilla/5.0"})
    response.raise_for_status()
    response.encoding = "utf-8"
    return response.text


def parse_feed(html):
    """Parse the table-style api_events markup into sorted event dicts.

    Malformed rows are skipped rather than aborting the whole run.
    """
    soup = BeautifulSoup(html, "html.parser")
    events = []
    for table in soup.find_all("table", class_="s-lc-ea-tb"):
        try:
            title_row = table.find("tr", class_="s-lc-ea-ttit")
            from_row = table.find("tr", class_="s-lc-ea-from")
            cat_row = table.find("tr", class_="s-lc-ea-tcat")
            if not (title_row and from_row):
                continue
            link = title_row.find("a")
            if not link:
                continue

            cells = from_row.find_all("td")
            if len(cells) < 2:
                continue
            when_text = cells[1].get_text(strip=True)
            # e.g. "11:00am Thursday, September 24, 2026"
            when = datetime.strptime(when_text, "%I:%M%p %A, %B %d, %Y")

            categories = []
            if cat_row:
                cat_cells = cat_row.find_all("td")
                if len(cat_cells) > 1:
                    categories = [c.strip() for c
                                  in cat_cells[1].get_text().split(",")
                                  if c.strip()]

            events.append({
                "title": link.get_text(strip=True),
                "link": link.get("href"),
                "when": when,
                "categories": categories,
            })
        except (ValueError, IndexError) as exc:
            print(f"  skipping unparseable event row: {exc}")
    events.sort(key=lambda e: e["when"])
    return events


def bucket_events(events):
    buckets = {box["header"]: [] for box in BOXES}
    for event in events:
        tokens = {c.casefold() for c in event["categories"]}
        for box in BOXES:
            if any(kw.casefold() in tokens for kw in box["keywords"]):
                buckets[box["header"]].append(event)
                break
        else:
            buckets[BOXES[-1]["header"]].append(event)
    return buckets


def event_label(event):
    return f"{event['when']:%B} {event['when'].day} - "


def event_lines(event):
    """Estimate rendered line count, memoized on the event dict."""
    if "_lines" not in event:
        event["_lines"] = max(
            1, -(-len(f"{event_label(event)}{event['title']}") // CHARS_PER_LINE))
    return event["_lines"]


def cap_buckets(buckets):
    """Fit sections within their line budgets, returning (capped, dropped)."""
    capped = {}
    dropped = []
    for header, events in buckets.items():
        budget = MAX_LINES.get(header)
        if budget is None:
            capped[header] = list(events)
            continue
        kept = []
        used = 0
        for event in events:
            lines = event_lines(event)
            if used + lines <= budget:
                kept.append(event)
                used += lines
            else:
                dropped.append((header, event))
        capped[header] = kept
    return capped, dropped


def split_dmds(events):
    """Split Do More events: first DMDS_MAIN_LINES in the main box, rest below."""
    main = []
    extra = []
    main_used = 0
    extra_used = 0
    for event in events:
        lines = event_lines(event)
        if main_used + lines <= DMDS_MAIN_LINES:
            main.append(event)
            main_used += lines
        elif extra_used + lines <= (MAX_LINES["Do More with Digital Scholarship"]
                                    - DMDS_MAIN_LINES):
            extra.append(event)
            extra_used += lines
    return main, extra


def term_title(events):
    when = events[0]["when"] if events else datetime.now()
    season = "Fall" if when.month >= 9 else "Summer" if when.month >= 6 else "Spring" if when.month >=3 else "Winter"
    return f"{season} {when.year}: Upcoming Workshops"


def find_shapes(slide):
    """Locate the title bar, the four quadrant boxes, and the overflow box."""
    title = None
    boxes = {}
    overflow = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.name == "Title 4" and shape.top < 1_000_000:
            title = shape
        elif shape.name == "Subtitle 2":
            header = shape.text_frame.paragraphs[0].text.strip()
            if header in BOX_HEADERS:
                boxes[header] = shape
        elif shape.name == "TextBox 2":
            overflow.append(shape)
    return title, boxes, overflow


def harvest_template(slide):
    """Deep-copy one existing event paragraph (date prefix + title) as template."""
    for shape in slide.shapes:
        if not shape.has_text_frame or shape.name != "Subtitle 2":
            continue
        for para in shape.text_frame.paragraphs:
            runs = para.runs
            if (len(runs) >= 2 and runs[1].text.strip()
                    and runs[0].text.strip().endswith(("-", ":"))):
                return copy.deepcopy(para._p)
    raise RuntimeError("No event-line paragraph found to use as a template")


def tighten_spacing(para, points=5):
    """Set space-after below the template's 6pt; schema-safe insertion."""
    spc_aft = para._p.get_or_add_pPr().get_or_add_spcAft()
    spc_pts = spc_aft.find(qn("a:spcPts"))
    if spc_pts is None:
        spc_pts = spc_aft.makeelement(qn("a:spcPts"), {})
        spc_aft.append(spc_pts)
        spc_pts.set("val", str(points * 100))
        return
    spc_pts.set("val", str(max(0, int(spc_pts.get("val")) - 100)))


def fill_box(shape, events, template, color, keep_header=True):
    tf = shape.text_frame
    first_to_remove = 1 if keep_header else 0
    for para in list(tf.paragraphs)[first_to_remove:]:
        para._p.getparent().remove(para._p)
    for _ in events:
        tf._txBody.append(copy.deepcopy(template))
    for para, event in zip(list(tf.paragraphs)[first_to_remove:], events):
        runs = para.runs
        runs[0].text = f"{event['when']:%B} {event['when'].day} - "
        runs[0].font.color.rgb = color
        runs[1].text = event["title"]
        runs[1].font.color.rgb = color
        runs[1].font.bold = False
        tighten_spacing(para)
        for extra in runs[2:]:
            extra._r.getparent().remove(extra._r)


def output_path():
    stamp = f"{datetime.now():%Y-%m-%d}"
    candidate = Path(f"SCDS Slide - {stamp}.pptx")
    version = 1
    while candidate.exists():
        version += 1
        candidate = Path(f"SCDS Slide - {stamp} v{version}.pptx")
    return candidate


def build_slide(events, deck_path=DECK_PATH):
    out_path = output_path()
    out_path.write_bytes(deck_path.read_bytes())
    prs = Presentation(str(out_path))

    # Older decks had two near-duplicate slides; keep only the second.
    sld_id_lst = prs.slides._sldIdLst
    if len(prs.slides) > 1:
        first = list(sld_id_lst)[0]
        prs.part.drop_rel(first.get(qn("r:id")))
        sld_id_lst.remove(first)

    slide = prs.slides[0]
    title_shape, boxes, overflow = find_shapes(slide)
    template = harvest_template(slide)

    title_shape.text_frame.paragraphs[0].runs[0].text = term_title(events)

    buckets, dropped = cap_buckets(bucket_events(events))
    for header, event in dropped:
        print(f"  [{header}] over limit, left off slide: "
              f"{event['when']:%b %d}  {event['title']}")

    dmds_main, dmds_extra = split_dmds(
        buckets["Do More with Digital Scholarship"])
    usage = {}
    for box in BOXES:
        box_events = buckets[box["header"]]
        if box["header"] == "Do More with Digital Scholarship":
            box_events = dmds_main
        shape = boxes.get(box["header"])
        if shape is not None:
            fill_box(shape, box_events, template, box["color"])
        usage[box["header"]] = sum(event_lines(e) for e in box_events)

    if overflow:
        fill_box(overflow[0], dmds_extra, template, WHITE, keep_header=False)
        usage["Do More with Digital Scholarship"] += sum(
            event_lines(e) for e in dmds_extra)
    for extra_box in overflow[1:]:
        extra_box.text_frame.clear()
    for header, used_lines in usage.items():
        print(f"  [{header}] {used_lines}/{MAX_LINES.get(header, '?')} lines")

    prs.save(str(out_path))
    return out_path


if __name__ == "__main__":
    html = fetch_feed()
    found = parse_feed(html)
    print(f"{len(found)} upcoming event(s) found")
    for event in found:
        print(f"  {event['when']:%b %d}  [{'; '.join(event['categories'])}]  {event['title']}")
    print(f"Wrote {build_slide(found)}")
